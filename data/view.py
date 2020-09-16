import argparse

parser = argparse.ArgumentParser()
setup_config = parser.add_argument_group('dataset setting')
setup_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp101')
setup_config.add_argument('--npy_name', type=str, dest='npy_name', default='trval')
setup_config.add_argument('--batch_size', type=int, dest='batch_size', default=60)
parser.print_help()
config, unparsed = parser.parse_known_args()

import tensorflow as tf
import numpy as np
import matplotlib.pyplot as plt
import pptx
from pptx.util import Inches
import sys, os
sys.path.append('/workspace/PycharmProjects/Rotator_loc')

from data.tf_data import DATA_PATH, data_setting, IMAGE_SIZE1


def data_loader(data_set):
    check_path = os.path.join(DATA_PATH, config.exp_name, 'view')
    if not os.path.exists(check_path):
        os.makedirs(check_path)

    with tf.Session() as sess:
        sess.run(data_set.val.init_op)
        num_examples, next_batch = data_set.val.data_length, data_set.val.next_batch

        count = 0
        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        print('num_iter: ', num_iter)

        prs = pptx.Presentation()
        prs.slide_width = Inches(10* 2)
        prs.slide_height = Inches(6 * 2)

        while True:
            try:
                 img1, xy_r1, xy_rp1, img1_t, lbl, name= sess.run(data_set.val.next_batch)

                 xr1 = IMAGE_SIZE1[1]*xy_r1[:,0].astype(np.int64)
                 yr1 = IMAGE_SIZE1[0]*xy_r1[:,1].astype(np.int64)

                 blank_slide_layout = prs.slide_layouts[6]
                 slide = prs.slides.add_slide(blank_slide_layout)

                 show_images(img1, name, xr1, yr1, count, num_rows=6, num_cols=10)
                 fig_name = '_'.join([config.exp_name, config.npy_name,
                                      '%03d' % count]) + '.png'
                 fig_path = os.path.join(check_path, fig_name)
                 plt.savefig(fig_path, bbox_inches='tight')
                 slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(10*2))
                 os.remove(fig_path)

                 count += 1

                 if count & 10 == 0:
                     print(count)
            except tf.errors.OutOfRangeError:
                break
    ppt_name = os.path.join(check_path, '_'.join([config.exp_name, config.npy_name])+'.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_images(images, names, x_coord, y_coord, idx, num_rows=5, num_cols=8, figsize=(10*2, 6*2)):
    plt.figure(figsize=figsize)
    num_figs = images.shape[0]  # num_rows * num_cols

    for j in range(num_figs):
        plt.subplot(num_rows, num_cols, j + 1)
        plt.imshow(np.squeeze(images[j]), cmap='gray')
        plt.annotate('O', xy=(x_coord[j], y_coord[j]), xytext=(x_coord[j], y_coord[j]), color='orange')
        plt.axis('off')
        plt.title(str(names[j].decode('utf-8')), fontsize=8, color='blue')


if __name__ == '__main__':

    data_set = data_setting(npy_name=config.npy_name)

    data_loader(data_set)