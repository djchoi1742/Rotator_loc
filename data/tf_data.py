import tensorflow as tf
import numpy as np
import pandas as pd
import os
import sys
import pydicom as dcm
import skimage.transform, scipy.misc
import re
import warnings
import argparse
import functools
import matplotlib.pyplot as plt
import pptx
from pptx.util import Inches

warnings.filterwarnings('ignore')

DATA_PATH = '/data/SNUBH/Rotator/ro_loc/'
RAW_PATH = '/data/SNUBH/Rotator/RAW/'
EXP_PATH = '/data/SNUBH/Rotator/ro_loc/'

IMAGE_SIZE1 = [256, 256]  # exp101: [512, 512]
IMAGE_SIZE2 = [512, 512]  # exp101: [256, 256]
CROP_RADIUS = 50

prop_h, prop_w = 0.5, 0.5

parser = argparse.ArgumentParser()
setup_config = parser.add_argument_group('data setting')
setup_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp104', help='experiment name')
setup_config.add_argument('--data_name', type=str, dest='data_name', default='trval', help='data name')
setup_config.add_argument('--batch_size', type=int, dest='batch_size', default=60, help='batch_size')
setup_config.add_argument('--train_csv', type=str, dest='train_csv', default=DATA_PATH + 'xlsx/exp308_train.xlsx')
setup_config.add_argument('--val_csv', type=str, dest='val_csv', default=DATA_PATH + 'xlsx/exp308_val.xlsx')
setup_config.add_argument('--view_type', type=int, dest='view_type', default=1)
config, unparsed = parser.parse_known_args()


def str_extract(string, text):
    extract = re.search(string, text)
    if extract is None:
        matching = None
    else:
        matching = extract.group()
    return matching


def data_setting(npy_name):
    only_val = config.train_csv == config.val_csv
    npy_path = os.path.join(DATA_PATH, config.exp_name, 'npy')
    data_dir = os.path.join(npy_path, npy_name)+'.npy'
    csv_path = os.path.join(DATA_PATH, config.exp_name, 'csv')
    if not os.path.exists(csv_path):
        os.makedirs(csv_path)

    csv = pd.DataFrame()
    csv_val = pd.read_excel(config.val_csv)
    csv_val['DATA_TYPE'] = 'val'
    csv = csv.append(csv_val, ignore_index=True)

    if only_val is False:
        csv_train = pd.read_excel(config.train_csv)
        csv_train['DATA_TYPE'] = 'train'
        csv = csv.append(csv_train, ignore_index=True)

    select_col = ['NUMBER', 'FOLDER_NAME', 'SIDE', 'LABEL_SST_BIN0',
                  'VIEW1', 'COORDS1_X', 'COORDS1_Y', 'SPACING1_X', 'SPACING1_Y',
                  'VIEW3', 'COORDS3_X', 'COORDS3_Y', 'SPACING3_X', 'SPACING3_Y',
                  'VIEW4', 'COORDS4_X', 'COORDS4_Y', 'SPACING4_X', 'SPACING4_Y',
                  'PATIENT_AGE', 'F', 'M', 'VAS_MED', 'TRAUMA0', 'TRAUMA1', 'TRAUMA2',
                  'DOMINANT0', 'DOMINANT1', 'DOMINANT2', 'DATA_TYPE'
                  ]

    csv = csv[select_col]

    def add_value(csv, view_type):
        view_name = 'VIEW'+view_type
        csv['FILES'+view_type] = csv.apply(
            lambda row: os.path.join(RAW_PATH, row['FOLDER_NAME'], row[view_name]), axis=1)
        csv[view_name + '_X'] = csv.apply(lambda row: row['COORDS'+view_type+'_X'], axis=1)
        csv[view_name + '_Y'] = csv.apply(lambda row: row['COORDS'+view_type+'_Y'], axis=1)
        csv[view_name + 'SPACING_X'] = csv.apply(lambda row: row['SPACING' + view_type + '_X'], axis=1)
        csv[view_name + 'SPACING_Y'] = csv.apply(lambda row: row['SPACING' + view_type + '_Y'], axis=1)
        csv['LABELS' + view_type] = csv.apply(lambda row: [row[view_name + '_X']] +
                                                          [row[view_name + '_Y']] +
                                                          [row[view_name + 'SPACING_X']] +
                                                          [row[view_name + 'SPACING_Y']] +
                                                          [row['LABEL_SST_BIN0']], axis=1)
        return csv

    csv = add_value(csv, '1')
    csv = add_value(csv, '3')
    csv = add_value(csv, '4')

    def files_labels_view(csv, view_type, data_type):
        files_view = csv[csv['DATA_TYPE'] == data_type]['FILES'+view_type].values
        labels_view = csv[csv['DATA_TYPE'] == data_type]['LABELS' + view_type].values
        return files_view, labels_view

    def info_csv(csv, data_type):
        file1, label1 = files_labels_view(csv, '1', data_type)
        file3, label3 = files_labels_view(csv, '3', data_type)
        file4, label4 = files_labels_view(csv, '4', data_type)
        id_data = csv[csv['DATA_TYPE'] == data_type]['NUMBER'].values

        data_info = pd.DataFrame({'FILES1': pd.Series(file1), 'LABELS1': pd.Series(label1),
                                  'FILES3': pd.Series(file3), 'LABELS3': pd.Series(label3),
                                  'FILES4': pd.Series(file4), 'LABELS4': pd.Series(label4),
                                  'ID': pd.Series(id_data)})
        data_info.to_csv(os.path.join(csv_path, '_'.join([npy_name, data_type+'.csv'])))
        return data_info

    val_csv = info_csv(csv, 'val')

    if only_val:
        data_set = DataSetting(data_dir=data_dir, only_val=only_val, batch_size=config.batch_size,
                               val_file=val_csv)
    else:
        train_csv = info_csv(csv, 'train')
        data_set = DataSetting(data_dir=data_dir, only_val=only_val, batch_size=config.batch_size,
                               val_file=val_csv, train_file=train_csv)
    return data_set


class DataSetting:
    def __init__(self, data_dir, only_val, batch_size, pos_aug=False, **kwargs):
        if not os.path.exists(data_dir):
            data_root = os.path.dirname(data_dir)
            if not os.path.exists(data_root):
                os.makedirs(data_root)

            if 'val_file' in kwargs:
                val_csv = kwargs['val_file']
                val_x1, val_y1 = val_csv['FILES1'].values, val_csv['LABELS1'].values
                val_x3, val_y3 = val_csv['FILES3'].values, val_csv['LABELS3'].values
                val_x4, val_y4 = val_csv['FILES4'].values, val_csv['LABELS4'].values
                val_id = val_csv['ID'].values
                if only_val:
                    np.save(data_dir, {'val_x1': val_x1, 'val_y1': val_y1,
                                       'val_x3': val_x3, 'val_y3': val_y3,
                                       'val_x4': val_x4, 'val_y4': val_y4, 'val_id': val_id,
                                       })
                else:
                    if 'train_file' in kwargs:
                        train_csv = kwargs['train_file']
                        train_x1, train_y1 = train_csv['FILES1'].values, train_csv['LABELS1'].values
                        train_x3, train_y3 = train_csv['FILES3'].values, train_csv['LABELS3'].values
                        train_x4, train_y4 = train_csv['FILES4'].values, train_csv['LABELS4'].values
                        train_id = train_csv['ID'].values

                        np.save(data_dir, {'train_x1': train_x1, 'train_y1': train_y1,
                                           'train_x3': train_x3, 'train_y3': train_y3,
                                           'train_x4': train_x4, 'train_y4': train_y4, 'train_id': train_id,
                                           'val_x1': val_x1, 'val_y1': val_y1,
                                           'val_x3': val_x3, 'val_y3': val_y3,
                                           'val_x4': val_x4, 'val_y4': val_y4, 'val_id': val_id,
                                           })
                    else:
                        raise AssertionError('training files or labels must be provided. please check npy file.')
            else:
                raise AssertionError('validation files or labels must be provided. please check npy file.')

        else:
            pre_built = np.load(data_dir).item()

            val_x1, val_y1 = pre_built['val_x1'], pre_built['val_y1']
            val_x3, val_y3 = pre_built['val_x3'], pre_built['val_y3']
            val_x4, val_y4 = pre_built['val_x4'], pre_built['val_y4']
            val_id = pre_built['val_id']

            self.data_length = len(val_id)
            self.val = self.RegSetting((val_x1, val_y1, val_x3, val_y3, val_x4, val_y4, val_id),
                                       batch_size=batch_size, shuffle=False)

            if only_val is False:

                train_x1, train_y1 = pre_built['train_x1'], pre_built['train_y1']
                train_x3, train_y3 = pre_built['train_x3'], pre_built['train_y3']
                train_x4, train_y4 = pre_built['train_x4'], pre_built['train_y4']
                train_id = pre_built['train_id']
                self.data_length = len(train_id) + len(val_id)

                if pos_aug:

                    index = np.asarray([v[-1] for v in train_y1])

                    train_x1 = np.concatenate([train_x1, train_x1[index == 1]], axis=0)
                    train_x3 = np.concatenate([train_x3, train_x3[index == 1]], axis=0)
                    train_x4 = np.concatenate([train_x4, train_x4[index == 1]], axis=0)
                    train_y1 = np.concatenate([train_y1, train_y1[index == 1]], axis=0)
                    train_y3 = np.concatenate([train_y3, train_y3[index == 1]], axis=0)
                    train_y4 = np.concatenate([train_y4, train_y4[index == 1]], axis=0)
                    train_id = np.concatenate([train_id, train_id[index == 1]], axis=0)

                np.random.seed(20201030)
                p = np.random.permutation(len(train_x1))
                train_x1, train_y1 = train_x1[p], train_y1[p]
                train_x3, train_y3 = train_x3[p], train_y3[p]
                train_x4, train_y4 = train_x4[p], train_y4[p]
                train_id = train_id[p]
                print(len(train_id))

                self.train = self.RegSetting((train_x1, train_y1, train_x3, train_y3, train_x4, train_y4, train_id),
                                              batch_size=batch_size, shuffle=True)

    class RegSetting:
        def __init__(self, files_n_labels, num_epochs=1, batch_size=1, shuffle=False, augmentation=False):
            self.file1, self.label1, self.file2, self.label2, self.file3, self.label3, self.id = files_n_labels
            self.data_length = len(self.id)

            self.neg_length = len([v[-1] for v in self.label1 if v[-1] == 0])
            self.pos_length = self.data_length - self.neg_length

            data_set = tf.data.Dataset.from_tensor_slices(tensors=
                                                         (self.file1, [v for v in self.label1],
                                                          self.file2, [v for v in self.label2],
                                                          self.file3, [v for v in self.label3],
                                                          [v for v in self.id],
                                                          ))
            if shuffle:
                data_set = data_set.shuffle(buffer_size=batch_size*100, reshuffle_each_iteration=True)

            def dcm_read_by_ftn(file1, label1, file2, label2, file3, label3, id):
                def each_read(file, label):
                    dcm_info = dcm.read_file(file.decode())
                    x, y, sx, sy = label[:-1]

                    img = dcm_info.pixel_array
                    h, w = dcm_info.Rows, dcm_info.Columns
                    h, w = np.int64(h), np.int64(w)
                    if str(dcm_info[0x28, 0x04].value) == 'MONOCHROME1':
                        white_img = np.full_like(img, np.max(img), img.dtype)
                        img = np.subtract(white_img, img)

                    center_x, center_y = int(w*prop_w), int(h*prop_h)

                    if h > w:
                        range_h1, range_h2 = center_y - int(w/2), center_y + int(w/2)
                        crop_img = img[range_h1:range_h2,:]
                        x_c, y_c = x, y - round(h/2) + round(w/2)

                    elif w > h:
                        range_w1, range_w2 = center_x - int(h/2), center_x + int(h/2)
                        crop_img = img[:,range_w1:range_w2]
                        x_c, y_c = x - round(w/2) + round(h/2), y

                    else:
                        crop_img = img
                        x_c, y_c = x, y

                    resize_img = np.expand_dims(skimage.transform.resize(crop_img, IMAGE_SIZE1,
                                                                         preserve_range=True), axis=-1)
                    x_r = round(x_c * (IMAGE_SIZE1[1] / crop_img.shape[1]))
                    y_r = round(y_c * (IMAGE_SIZE1[0] / crop_img.shape[0]))

                    x_r, y_r = np.int64(x_r), np.int64(y_r)

                    x_rp, y_rp = np.float32(x_r / resize_img.shape[1]), np.float32(y_r / resize_img.shape[0])

                    resize_img = resize_img.astype(np.float32)
                    resize_img = (resize_img - np.mean(resize_img)) / np.std(resize_img)

                    return resize_img, np.array([x_rp, y_rp]), file.decode(), np.int64(label[-1])

                name = id.decode()

                img1, xy_rp1, file1, lbl = each_read(file1, label1)
                img2, xy_rp2, file2, _ = each_read(file2, label2)
                img3, xy_rp3, file3, _ = each_read(file3, label3)

                return img1, xy_rp1, file1, img2, xy_rp2, file2, img3, xy_rp3, file3, lbl, name

            data_set = data_set.map(num_parallel_calls=2,
                                    map_func=lambda file1, label1, file2, label2, file3, label3, id:
                                    tuple(tf.py_func(func=dcm_read_by_ftn,
                                                     inp=[file1, label1, file2, label2, file3, label3, id],
                                                     Tout=[tf.float32, tf.float32, tf.string,
                                                           tf.float32, tf.float32, tf.string,
                                                           tf.float32, tf.float32, tf.string,
                                                           tf.int64, tf.string]
                                                     )))
            if num_epochs == 0:
                data_set = data_set.repeat(count=num_epochs)  # raise out-of-range error when num_epochs done
                data_set = data_set.batch(batch_size=batch_size, drop_remainder=True)
            else:
                data_set = data_set.batch(batch_size)
                data_set = data_set.repeat(count=num_epochs)
            iterator = data_set.make_initializable_iterator()

            self.init_op = iterator.initializer
            self.next_batch = iterator.get_next()


def calculate_crop_coord(center_x, center_y, spacing_x, spacing_y, radius):
    x1, y1 = int(center_x - radius / spacing_x), int(center_y - radius / spacing_y)
    x2, y2 = int(center_x + radius / spacing_x), int(center_y + radius / spacing_y)
    return [x1, y1, x2, y2]


def local_process(value):
    if value < 0:
        local_value = 0.0
    elif value > IMAGE_SIZE1[0]:
        local_value = float(IMAGE_SIZE1[0])
    else:
        local_value = value
    return float(local_value)


class LocalizeResult:
    def __init__(self, files, locals):
        files = files.astype(str)
        self.loc = list(map(self.loc_result, files.tolist(), locals.tolist()))

    def loc_result(self, file, local):
        read_file = dcm.read_file(file)
        image = read_file.pixel_array
        h, w = image.shape

        x_r, y_r = local_process(local[0]), local_process(local[1])

        if h > w:
            x_c, y_c = (w/IMAGE_SIZE1[1])*x_r, (w/IMAGE_SIZE1[1])*y_r
            x, y = round(x_c), round(y_c + h/2 - w/2)
        elif w > h:
            x_c, y_c = (h/IMAGE_SIZE1[0])*x_r, (h/IMAGE_SIZE1[1])*y_r
            x, y = round(x_c + w/2 - h/2), round(y_c)
        else:
            x, y = (w/IMAGE_SIZE1[1])*x_r, (w/IMAGE_SIZE1[1])*y_r

        return [x, y]


class ImageProcess:
    def __init__(self, files, locals, augmentation):
        files = files.astype(str)
        self.crops = np.asarray(list(map(functools.partial(self.dcm_extract, augmentation=augmentation),
                                         files.tolist(), locals.tolist())))

    def dcm_extract(self, file, local, augmentation):
        read_file = dcm.read_file(file)
        try:
            spacing_x, spacing_y = read_file.ImagerPixelSpacing[1], read_file.ImagerPixelSpacing[0]
        except:
            spacing_x, spacing_y = read_file.PixelSpacing[1], read_file.PixelSpacing[0]
        image = read_file.pixel_array
        h, w = image.shape

        x_r, y_r = local_process(local[0]), local_process(local[1])

        if h > w:
            x_c, y_c = (w/IMAGE_SIZE1[1])*x_r, (w/IMAGE_SIZE1[1])*y_r
            x, y = round(x_c), round(y_c + h/2 - w/2)
        elif w > h:
            x_c, y_c = (h/IMAGE_SIZE1[0])*x_r, (h/IMAGE_SIZE1[1])*y_r
            x, y = round(x_c + w/2 - h/2), round(y_c)
        else:
            x, y = (w/IMAGE_SIZE1[1])*x_r, (w/IMAGE_SIZE1[1])*y_r

        x_range1, x_range2 = int(x - CROP_RADIUS / spacing_x), int(x + CROP_RADIUS / spacing_x)
        y_range1, y_range2 = int(y - CROP_RADIUS / spacing_y), int(y + CROP_RADIUS / spacing_y)

        if augmentation:
            shift_x, shift_y = np.random.randint(w // 10), np.random.randint(h // 10)
            shift_x = -shift_x if np.random.rand() <= 0.5 else shift_x
            shift_y = -shift_y if np.random.rand() <= 0.5 else shift_y

            x_range1, x_range2 = x_range1 - shift_x, x_range2 - shift_x
            y_range1, y_range2 = y_range1 - shift_y, y_range2 - shift_y

        crop_img = image[max(0, y_range1):min(h, y_range2), max(0, x_range1):min(w, x_range2)]
        crop_img = np.expand_dims(skimage.transform.resize(crop_img, IMAGE_SIZE2, preserve_range=True), axis=-1)

        if augmentation and np.random.randint(2) == 1:
            crop_img = np.fliplr(crop_img)

        crop_img = (crop_img - np.mean(crop_img)) / np.std(crop_img)
        return crop_img


def select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3):
    if config.view_type == 1:
        images, xy_rp, files = images1, xy_rp1, files1
    elif config.view_type == 3:
        images, xy_rp, files = images2, xy_rp2, files2
    elif config.view_type == 4:
        images, xy_rp, files = images3, xy_rp3, files3
    else:
        raise ValueError('Error! Invalid view type.')
    return images, xy_rp, files


def test_data_loader(data_set):
    check_path = os.path.join(EXP_PATH, config.exp_name, 'view')
    if not os.path.exists(check_path):
        os.makedirs(check_path)

    with tf.Session() as sess:
        sess.run(data_set.val.init_op)
        num_examples, next_batch = data_set.val.data_length, data_set.val.next_batch

        count = 0
        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        print('num_iter: ', num_iter)

        prs = pptx.Presentation()
        prs.slide_width = Inches(10 * 2)
        prs.slide_height = Inches(6 * 2)

        while count < num_iter:
            images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, \
            labels, names = sess.run(data_set.val.next_batch)

            images, xy_rp, files = \
                select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            show_images(images, names, num_rows=6, num_cols=10, fig_size=(10*2, 6*2))
            fig_name = '_'.join([config.exp_name, config.data_name, '%03d' % count]) + '.png'

            fig_path = os.path.join(check_path, fig_name)
            plt.savefig(fig_path, bbox_inches='tight')

            slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(10 * 2))
            os.remove(fig_path)
            count += 1

            if count % 10 == 0:
                print(count)

    ppt_name = os.path.join(check_path, '_'.join([config.exp_name, config.data_name]) + '.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_images(images, names, num_rows=6, num_cols=10, fig_size=(10*2, 6*2)):
    plt.figure(figsize=fig_size)
    num_figs = images.shape[0]  # num_rows * num_cols
    # num_chars = 5  # num of chars to show in names

    for j in range(num_figs):
        plt.subplot(num_rows, num_cols, j + 1)
        plt.imshow(np.squeeze(images[j]), cmap='gray')
        plt.axis('off')
        img_name = os.path.basename(names[j])

        plt_name = '_'.join([str(img_name.decode('utf-8'))])
        plt.title(plt_name, fontsize=8, color='blue')


if __name__ == '__main__':
    os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'
    os.environ['CUDA_VISIBLE_DEVICES'] = '3'

    d_set = data_setting(npy_name=config.data_name)
    test_data_loader(d_set)

