import argparse

license = """
Copyright ⓒ Dongjun Choi, Kyong Joon Lee
Department of Radiology at Seoul National University Bundang Hospital. \n
If you have any question, please email us for assistance: dongyul.oh@snu.ac.kr \n """
parser = argparse.ArgumentParser(formatter_class=argparse.RawDescriptionHelpFormatter, \
                                 description='', epilog=license, add_help=False)

network_config = parser.add_argument_group('network setting (must be provided)')

network_config.add_argument('--data_path', type=str, dest='data_path', default='/data/SNUBH/Rotator/ro_loc/')
network_config.add_argument('--exp_name', type=str, dest='exp_name', default='exp103')
network_config.add_argument('--model1_name', type=str, dest='model1_name', default='Model03')
network_config.add_argument('--model2_name', type=str, dest='model2_name', default='Model04')
network_config.add_argument('--train', type=lambda x: x.title() in str(True), dest='train', default=False)
network_config.add_argument('--batch_size', type=int, dest='batch_size', default=8)
network_config.add_argument('--num_epoch', type=int, dest='num_epoch', default=100)  # only use training
network_config.add_argument('--trial_serial', type=int, dest='trial_serial', default=7)
network_config.add_argument('--npy_name', type=str, dest='npy_name', default='trval.npy')
network_config.add_argument('--max_keep', type=int, dest='max_keep', default=5)  # only use training
network_config.add_argument('--num_weight', type=int, dest='num_weight', default=5)  # only use validation
network_config.add_argument('--block_rep1', type=str, dest='block_rep1', default='2,2,2,2,2')
network_config.add_argument('--k_p1', type=str, dest='k_p1', default='1,1,1,1,1')
network_config.add_argument('--block_rep2', type=str, dest='block_rep2', default='3,3,3,3')
network_config.add_argument('--k_p2', type=str, dest='k_p2', default='1,1,1,1')
network_config.add_argument('--use_se', type=lambda x: x.title() in str(True), dest='use_se', default=False)
network_config.add_argument('--learning_rate', type=float, dest='learning_rate', default=0.01)
network_config.add_argument('--decay_steps', type=int, dest='decay_steps', default=7500)  # prev: 7500
network_config.add_argument('--decay_rate', type=float, dest='decay_rate', default=0.94)
network_config.add_argument('--decay', type=float, dest='decay', default=0.9)
network_config.add_argument('--alpha', type=float, dest='alpha', default=0.3)  # only use training
network_config.add_argument('--gamma', type=float, dest='gamma', default=2)  # only use training
network_config.add_argument('--pos_aug', type=lambda x: x.title() in str(True), dest='pos_aug', default=False)
network_config.add_argument('--view_type', type=int, dest='view_type', default=1)

#parser.print_help()
config, unparsed = parser.parse_known_args()

import sys, os
sys.path.append('/home/chzze/bitbucket/Rotator_loc')

os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'  # ignore SSE instruction warning on tensorflow

import tensorflow as tf
import numpy as np
import sklearn.metrics  # roc curve
import matplotlib.pyplot as plt
import pandas as pd
import json, pptx
from pptx.util import Inches

trial_serial_str = '%03d' % config.trial_serial
log_path = os.path.join(config.data_path, config.exp_name, config.model1_name, 'logs-%s' % trial_serial_str)
ppt_path = os.path.join(config.data_path, config.exp_name, config.model1_name, 'ppt-%s' % trial_serial_str)
result_path = os.path.join(config.data_path, config.exp_name, config.model1_name, 'result-%s' % trial_serial_str)
png_path = os.path.join(ppt_path, 'png')
ckpt_path = os.path.join(result_path, 'ckpt')
npy_path = os.path.join(config.data_path, config.exp_name, 'npy')


if not os.path.exists(result_path):
    os.makedirs(result_path)

if not os.path.exists(ppt_path):
    os.makedirs(ppt_path)

if not os.path.exists(png_path):
    os.makedirs(png_path)


from data.tf_data import DataSetting, ImageProcess, LocalizeResult
import models.model_b as model_b
from tf_utils.tensor_board import TensorBoard
print('/'.join([npy_path, config.npy_name]))
data_set = DataSetting(data_dir=os.path.join(npy_path, config.npy_name), batch_size=config.batch_size,
                       only_val=bool(1 - config.train), pos_aug=config.pos_aug)

if config.train:
    block_rep1, k_p1 = config.block_rep1, config.k_p1
    block_rep2, k_p2 = config.block_rep2, config.k_p2
    view_type = config.view_type
else:
    block_rep1, k_p1 = model_b.get_info(result_path, 'BLOCK_REP1'), model_b.get_info(result_path, 'K_P1')
    block_rep2, k_p2 = model_b.get_info(result_path, 'BLOCK_REP2'), model_b.get_info(result_path, 'K_P2')
    view_type = model_b.get_info(result_path, 'VIEW_TYPE')

infer1_name = 'Inference' + config.model1_name
infer2_name = 'Inference' + config.model2_name
model1 = getattr(model_b, infer1_name)(block_rep=block_rep1, k_p=k_p1, view=view_type)
model2 = getattr(model_b, infer2_name)(trainable=config.train, view=view_type, block_rep=block_rep2, k_p=k_p2,
                                       use_se=config.use_se, alpha=config.alpha, gamma=config.gamma)



sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
sess_config.gpu_options.allow_growth = True
sess = tf.Session(config=sess_config)


def select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, view):
    if view == 1:
        images, xy_rp, files = images1, xy_rp1, files1
    elif view == 3:
        images, xy_rp, files = images2, xy_rp2, files2
    elif view == 4:
        images, xy_rp, files = images3, xy_rp3, files3
    else:
        raise ValueError('Error! Invalid view type.')
    return images, xy_rp, files


def training():
    info_log = {
        'NPY_NAME': config.npy_name,
        'BATCH_SIZE': config.batch_size,
        'LEARNING_RATE': config.learning_rate,
        'DECAY_STEPS': config.decay_steps,
        'DECAY_RATE': config.decay_rate,
        'NUM_EPOCH': config.num_epoch,
        'BLOCK_REP1': config.block_rep1,
        'K_P1': config.k_p1,
        'BLOCK_REP2': config.block_rep2,
        'K_P2': config.k_p2,
        'USE_SE': config.use_se,
        'POS_AUG': config.pos_aug,
        'FOCAL_LOSS_ALPHA': config.alpha,
        'FOCAL_LOSS_GAMMA': config.gamma,
        'VIEW_TYPE': config.view_type
    }

    with open(os.path.join(result_path, '.info'), 'w') as f:
        f.write(json.dumps(info_log, indent=4, sort_keys=True))
        f.close()

    ts_board = TensorBoard(log_dir=log_path, overwrite=True)
    mse = tf.get_variable(name='MSE', shape=[], trainable=False,
                          initializer=tf.zeros_initializer(), collections=['scalar'])
    loss = tf.get_variable(name='Entropy', shape=[], trainable=False,
                           initializer=tf.zeros_initializer(), collections=['scalar'])
    auc_rec = tf.get_variable(name='AUC', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                              collections=['scalar'])
    accuracy_rec = tf.get_variable(name='Accuracy', shape=[], trainable=False, initializer=tf.zeros_initializer(),
                                   collections=['scalar'])

    ts_board.init_scalar(collections=['scalar'])
    ts_board.init_images(collections=['images'], num_outputs=4)

    train_op = model_b.optimize(reg_loss=model1.mse, ce_loss=model2.loss, learning_rate=config.learning_rate,
                                decay_rate=config.decay_rate, decay_steps=config.decay_steps)

    sess.run([tf.global_variables_initializer(), tf.local_variables_initializer()])
    saver = tf.train.Saver(max_to_keep=config.max_keep)

    result_name = '_'.join([config.exp_name, config.model1_name, trial_serial_str])+'.csv'
    df_csv = pd.DataFrame({'WEIGHT_PATH': pd.Series(), 'AUC': pd.Series()})

    crt_step, crt_epoch = None, None
    perf_per_epoch, max_perf_per_epoch, max_crt_step = [], [], []

    try:
        while True:
            sess.run([data_set.train.init_op, data_set.val.init_op])
            train_loss_batch, train_mse_batch, train_acc_batch = [], [], []
            train_x, train_y = [], []

            train_length = data_set.train.data_length
            num_iter_train = int(np.ceil(float(train_length) / config.batch_size))
            train_step = 0

            feed_dict = {}
            while train_step < num_iter_train:
                images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, \
                labels, names = sess.run(data_set.train.next_batch)

                images, xy_rp, files = \
                    select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, view_type)

                feed_local = {model1.images: images, model1.locals: xy_rp, model1.is_training: True}
                train_mse, train_loc = sess.run([model1.mse, model1.predict], feed_dict=feed_local)

                train_process = ImageProcess(files=files, locals=train_loc, augmentation=True)

                feed_dict = {model2.images: train_process.crops, model2.labels: labels, model2.is_training: True,
                             model1.images: images, model1.locals: xy_rp, model1.is_training: True}

                _, train_loss, train_prob, train_acc, crt_step, crt_epoch = \
                    sess.run([train_op, model2.loss, model2.prob, model2.accuracy,
                              tf.train.get_global_step(), model2.global_epoch], feed_dict=feed_dict)

                sys.stdout.write('Step: {0:>4d} ({1})\r'.format(crt_step, int(crt_epoch)))

                train_x.extend(train_prob)
                train_y.extend(labels)
                train_acc_batch.append(train_acc)
                train_mse_batch.append(train_mse)
                train_loss_batch.append(train_loss)

                train_step += 1

            sess.run(tf.assign_add(model2.global_epoch, 1))

            fpr, tpr, _ = sklearn.metrics.roc_curve(train_y, train_x, drop_intermediate=False)
            train_auc = sklearn.metrics.auc(fpr, tpr)

            feed_dict.update({loss: np.mean(train_loss_batch), mse: np.mean(train_mse_batch),
                              auc_rec: train_auc, accuracy_rec: np.mean(train_acc_batch)})

            ts_board.add_summary(sess=sess, feed_dict=feed_dict, log_type='train')

            val_length = data_set.val.data_length
            num_iter_val = int(np.ceil(float(val_length) / config.batch_size))
            val_step = 0

            val_loss_batch, val_mse_batch, val_acc_batch = [], [], []
            val_x, val_y = [], []

            while val_step < num_iter_val:
                sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_loss_batch)+1, num_iter_val))

                images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, \
                labels, names = sess.run(data_set.val.next_batch)

                images, xy_rp, files = \
                    select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, view_type)

                feed_local = {model1.images: images, model1.locals: xy_rp, model1.is_training: False}
                val_mse, val_loc = sess.run([model1.mse, model1.predict], feed_dict=feed_local)

                val_process = ImageProcess(files=files, locals=val_loc, augmentation=False)

                feed_dict = {model2.images: val_process.crops, model2.labels: labels, model2.is_training: False,
                             model1.images: images, model1.locals: xy_rp, model1.is_training: False}

                val_loss, val_prob, val_acc = sess.run([model2.loss, model2.prob, model2.accuracy],
                                                       feed_dict=feed_dict)

                val_x.extend(val_prob)
                val_y.extend(labels)
                val_acc_batch.append(val_acc)
                val_mse_batch.append(val_mse)
                val_loss_batch.append(val_loss)

                val_step += 1

            fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, val_x, drop_intermediate=False)
            val_auc = sklearn.metrics.auc(fpr, tpr)

            feed_dict.update({loss: np.mean(val_loss_batch), mse: np.mean(val_mse_batch),
                              auc_rec: val_auc, accuracy_rec: np.mean(val_acc_batch)})

            ts_board.add_summary(sess=sess, feed_dict=feed_dict, log_type='val')
            ts_board.display_summary(time_stamp=True)

            crt_epoch += 1
            if crt_epoch % 1 == 0:
                perf_per_epoch.append(val_auc)

                if crt_epoch < config.max_keep + 1:
                    max_crt_step.append(crt_step)
                    max_perf_per_epoch.append(val_auc)

                    saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'),
                               global_step=crt_step)
                    df_csv.loc[crt_step, 'WEIGHT_PATH'] = os.path.join(log_path, 'model.ckpt-'+str(crt_step))
                    df_csv.loc[crt_step, 'AUC'] = val_auc

                elif val_auc > min(df_csv['AUC'].tolist()):
                    df_csv = df_csv.drop(max_crt_step[0])
                    max_crt_step.pop(0)
                    max_crt_step.append(crt_step)
                    max_perf_per_epoch.pop(0)
                    max_perf_per_epoch.append(val_auc)

                    saver.save(sess=sess, save_path=os.path.join(log_path, 'model.ckpt'),
                               global_step=crt_step)
                    df_csv.loc[crt_step, 'WEIGHT_PATH'] = os.path.join(log_path, 'model.ckpt-'+str(crt_step))
                    df_csv.loc[crt_step, 'AUC'] = val_auc

                df_csv.to_csv(os.path.join(result_path, result_name))

                if crt_epoch == config.num_epoch:
                    break
                    
        print('Training Complete...\n')
        sess.close()

    except KeyboardInterrupt:
        print('Result saved')
        df_csv.to_csv(os.path.join(result_path, result_name))


def validation():
    sess_config = tf.ConfigProto(log_device_placement=False, allow_soft_placement=True)
    sess_config.gpu_options.allow_growth = True

    saver = tf.train.Saver()
    num_examples = data_set.val.data_length

    weight_auc_path = os.path.join(config.data_path, config.exp_name, config.model1_name,
                                   'result-%03d' % config.trial_serial)
    weight_df_csv = pd.read_csv(os.path.join(weight_auc_path, '_'.join([config.exp_name, config.model1_name,
                                                                        '%03d' % config.trial_serial])+'.csv'))

    weight_df_csv = weight_df_csv.sort_values('AUC', ascending=False)
    all_ckpt_paths = list(weight_df_csv['WEIGHT_PATH'][0:int(config.num_weight)])

    num_ckpt = len(all_ckpt_paths)
    print('num_ckpt: ', num_ckpt)

    imgs = np.zeros([num_examples, model2.img_h, model2.img_w, model2.img_c])
    cams = np.zeros([num_ckpt, num_examples, model2.img_h, model2.img_w, model2.img_c])

    lbls = np.zeros([num_examples, ], dtype=np.int32)
    lgts = np.zeros([num_ckpt, num_examples, 1])
    probs = np.zeros([num_ckpt, num_examples, 1])
    p_locs = np.zeros([num_ckpt, num_examples, 2])
    mses_x = np.zeros([num_ckpt, num_examples, 1], dtype=np.float32)
    mses_y = np.zeros([num_ckpt, num_examples, 1], dtype=np.float32)

    val_x, val_y = None, None
    for ckpt_idx, ckpt_path in enumerate(all_ckpt_paths):
        w_path = os.path.basename(ckpt_path)
        ckpt_path = os.path.join(log_path, w_path)
        print('Restoring: '+ckpt_path)

        sess = tf.Session(config=sess_config)
        saver.restore(sess, ckpt_path)

        sess.run(data_set.val.init_op)

        val_x, val_y = [], []
        num_iter = int(np.ceil(float(num_examples) / config.batch_size))
        step = 0

        while step < num_iter:
            sys.stdout.write('Evaluation [{0}/{1}]\r'.format(len(val_y) // config.batch_size + 1,
                             -(-data_set.val.data_length // config.batch_size)))

            images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, \
            labels, names = sess.run(data_set.val.next_batch)

            images, xy_rp, files = \
                select_view(images1, xy_rp1, files1, images2, xy_rp2, files2, images3, xy_rp3, files3, view_type)

            feed_local = {model1.images: images, model1.locals: xy_rp, model1.is_training: False}
            val_mse, val_loc = sess.run([model1.mse, model1.predict], feed_dict=feed_local)

            val_process = ImageProcess(files=files, locals=val_loc, augmentation=False)

            feed_dict = {model2.images: val_process.crops, model2.labels: labels, model2.is_training: False,
                         model1.images: images, model1.locals: xy_rp, model1.is_training: False}

            val_res, val_mse, val_loc, val_lgt, val_prob, val_acc, val_cam = \
                sess.run([model1.residual, model1.mse, model1.predict, model2.logits, model2.prob,
                          model2.accuracy, model2.local], feed_dict=feed_dict)

            loc_result = LocalizeResult(files=files, locals=val_loc)
            loc_xy = loc_result.loc

            val_x.extend(val_prob)
            val_y.extend(labels)

            res_x, res_y = np.expand_dims(val_res[:, 0], axis=-1), np.expand_dims(val_res[:, 1], axis=-1)

            p_locs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = loc_xy
            mses_x[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = res_x
            mses_y[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = res_y

            lgts[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = val_lgt
            probs[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = val_prob
            cams[ckpt_idx, step * config.batch_size:step * config.batch_size + len(labels)] = val_cam

            if ckpt_idx == 0:
                lbls[step * config.batch_size:step * config.batch_size + len(labels)] = labels
                imgs[step * config.batch_size:step * config.batch_size + len(labels)] = val_process.crops
                if step == 0:
                    pd.DataFrame(np.squeeze(val_cam[0])).to_excel(os.path.join(result_path,
                                                                               str(config.view_type) + '_ex0.xlsx'))

            step += 1
        sess.close()

    p_locs = np.mean(p_locs, axis=0)
    lgts, probs, cams = np.mean(lgts, axis=0), np.mean(probs, axis=0), np.mean(cams, axis=0)
    lgts_1 = np.squeeze(np.array(lgts))
    prob_1 = np.squeeze(np.array(probs))

    result_csv = pd.DataFrame({'NUMBER': data_set.val.id, 'LOGIT': lgts_1, 'PROB': prob_1,
                               'LABEL': np.array(lbls), 'LOC_X': p_locs[:, 0], 'LOC_Y': p_locs[:, 1]})

    result_name = '_'.join([config.model1_name, config.npy_name, trial_serial_str, '%03d' % config.num_weight])+'.csv'

    result_csv.to_csv(os.path.join(result_path, result_name), index=False)

    mse_x, mse_y = np.mean(mses_x), np.mean(mses_y)

    fpr, tpr, _ = sklearn.metrics.roc_curve(val_y, prob_1, drop_intermediate=False)
    val_auc = sklearn.metrics.auc(fpr, tpr)
    print('Validation MSE: %04f, %04f' % (mse_x, mse_y))
    print('Validation AUC: ', val_auc)
    print('Validation Complete...\n')

    prs = pptx.Presentation()
    prs.slide_width, prs.slide_height = Inches(8*2), Inches(5*2)

    plt_batch = 20
    plt_step = 0
    plt_iter, plt_examples = int(np.ceil(num_examples / plt_batch)), num_examples

    while plt_step < plt_iter:
        if plt_examples >= plt_batch:
            len_batch = plt_batch
        else:
            len_batch = plt_examples

        labels_batch = lbls[plt_step * plt_batch:plt_step * plt_batch + len_batch]
        names_batch = data_set.val.id[plt_step * plt_batch:plt_step*plt_batch + len_batch]
        probs_batch = probs[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        images_batch = imgs[plt_step * plt_batch:plt_step*plt_batch + len_batch]
        cams_batch = cams[plt_step * plt_batch:plt_step * plt_batch + len_batch]

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        show_cam(cams_batch, probs_batch, images_batch, labels_batch, names_batch, 'LABEL')
        fig_name = '_'.join([config.exp_name, config.model1_name, config.npy_name, trial_serial_str,
                             '%03d' % plt_step])
        fig_path = os.path.join(png_path, fig_name)+'.png'
        plt.savefig(fig_path, bbox_inches='tight')
        slide.shapes.add_picture(fig_path, Inches(0), Inches(0), width=Inches(8*2))
        # os.remove(fig_path)
        plt_step += 1
        plt_examples -= plt_batch

    print('plt_examples check: ', plt_examples)
    ppt_name = os.path.join(ppt_path, '_'.join([config.exp_name, config.model1_name, config.npy_name,
                                                trial_serial_str, '%03d' % config.num_weight])+'.pptx')
    prs.save(ppt_name)
    print('Saved: ', ppt_name)


def show_cam(cams, probs, images, labels, names, side_label, num_rows=5, num_cols=8, figsize=(8*2, 5*2)):
    batch_size = cams.shape[0]
    fig, ax = plt.subplots(num_rows, num_cols, figsize=figsize)
    axoff_fun = np.vectorize(lambda ax: ax.axis('off'))
    axoff_fun(ax)

    for i in range(batch_size):
        prob, lbl = '%.2f' % probs[i], int(labels[i])
        show_image, cam = np.squeeze(images[i]), np.squeeze(cams[i])
        img_row, img_col = int(i % num_rows), int(i / num_rows) * 2

        ori_title = ' '.join([names[i], side_label, ': '+str(lbl)])
        cam_title = side_label+' Pred: '+str(prob)

        ax[img_row, img_col].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(show_image, cmap='bone')
        ax[img_row, img_col+1].imshow(cam, cmap=plt.cm.seismic, alpha=0.5, interpolation='nearest')

        if (lbl == 0 and probs[i] < 0.5) or (lbl == 1 and probs[i] >= 0.5):
            txt_color = 'blue'
        else:
            txt_color = 'red'
        ax[img_row, img_col].set_title(ori_title, fontsize=7, color=txt_color)
        ax[img_row, img_col+1].set_title(cam_title, fontsize=7, color=txt_color)


if __name__ == '__main__':
    if config.train:
        training()
    else:
        validation()
